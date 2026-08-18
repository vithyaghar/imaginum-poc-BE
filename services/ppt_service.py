#!/usr/bin/env python
# -- coding:utf-8 --
###
# Filename: c:\Projects\imaginum\backend\services\ppt_service.py
# Path: c:\Projects\imaginum\backend\services
# Created Date: Tuesday, March 10th 2026, 12:17:54 pm
# Author: Vithyaghar M
#
# Copyright (c) 2026 Trinom Digital
###
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import io
import base64
from typing import Optional
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate

# from langchain.agents import create_agent, AgentState
import json
from dotenv import load_dotenv

load_dotenv()
# ─── Helpers ────────────────────────────────────────────────────────────────


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


async def ai_parse_markdown_to_slides(markdown: str):

    print("=============> ai_parse_markdown_to_slides")

    prompt = ChatPromptTemplate.from_template(
        """
You are a presentation planning AI.

Convert the following Markdown campaign document into a structured slide deck.

CRITICAL RULES:
- Do NOT lose important details.
- Do NOT aggressively summarize.
- Preserve all key insights, numbers, strategies, and reasoning.
- Expand slightly where needed to make slides meaningful and clear.
- Avoid copying the entire markdown verbatim — adapt it into presentation format.

CONTENT GUIDELINES:
- Each slide should feel complete and informative on its own.
- Use bullets for clarity, but include enough detail (not 1–2 word bullets).
- Use the "body" field to:
  - Add context
  - Explain reasoning
  - Include supporting detail that doesn’t fit in bullets
- Important: If a section has rich content, distribute it across bullets + body.
- Do NOT compress multi-line ideas into a single vague sentence.

BULLET RULES:
- Each bullet should be 1–2 lines max
- Preserve key phrases (e.g., numbers, KPIs, budgets, timelines)
- Do NOT oversimplify (e.g., avoid "Increase sales" → keep specifics)

SLIDE STRUCTURE:
- First slide MUST be a title slide (`is_title_slide: true`)
- Convert each major section into 1 slide (or more if content is dense)
- Maintain logical grouping of ideas

Return ONLY valid JSON. No explanations.

Output format:

[
  {{
    "title": "string",
    "layout": "title|content",
    "bullets": ["point 1", "point 2"],
    "body": "supporting explanation or elaboration",
    "is_title_slide": true
  }}
]

Markdown:

{markdown}
"""
    )

    slide_llm = ChatGoogleGenerativeAI(
        model="models/gemini-3.5-flash",
        temperature=0,
        vertexai=False,
        max_output_tokens=4096,
        google_api_key=os.getenv("GOOGLE_API_KEY"),
    )

    chain = prompt | slide_llm

    result = await chain.ainvoke({"markdown": markdown})

    content = result.content

    print("LLM RAW OUTPUT TYPE:", type(content))
    print("LLM RAW OUTPUT:", content)

    # Gemini sometimes wraps the response like this:
    # [{'type': 'text', 'text': 'JSON_STRING'}]

    if isinstance(content, list) and "text" in content[0]:
        json_string = content[0]["text"]
        return json.loads(json_string)

    if isinstance(content, str):
        return json.loads(content)

    if isinstance(content, list):
        return content

    raise ValueError(f"Unexpected LLM output: {content}")


def normalize_bullets(slides):
    for slide in slides:
        bullets = slide.get("bullets", [])
        normalized = []

        for bullet in bullets:
            if isinstance(bullet, str):
                normalized.append({"text": bullet, "level": 0, "bold": False})
            else:
                normalized.append(bullet)

        slide["bullets"] = normalized

    return slides


# ─── Core Slide Builders ─────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def build_title_slide(prs: Presentation, slide_data: dict, palette: dict):

    print("=============> build_title_slide")

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb(palette["dark"])

    # Accent bar (left edge)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(palette["accent"])
    bar.line.fill.background()

    # Title text
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(10), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = slide_data["title"]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(palette["light"])
    run.font.name = palette["header_font"]

    # Subtitle / body
    if slide_data["body"].strip():
        txb2 = slide.shapes.add_textbox(
            Inches(0.6), Inches(4.5), Inches(10), Inches(1.2)
        )
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = slide_data["body"].strip()
        run2.font.size = Pt(18)
        run2.font.color.rgb = hex_to_rgb(palette["muted"])
        run2.font.name = palette["body_font"]

    return slide


def build_content_slide(
    prs: Presentation, slide_data: dict, palette: dict, slide_index: int
):

    print("=============> build_content_slide")

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Alternate light / slightly-off-white backgrounds for rhythm
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb(
        palette["light"] if slide_index % 2 == 0 else palette["light_alt"]
    )

    # Header bar
    header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = hex_to_rgb(palette["dark"])
    header_bar.line.fill.background()

    # Slide title in header
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(11), Inches(0.8))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = slide_data["title"]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(palette["light"])
    run.font.name = palette["header_font"]

    # Accent side stripe
    stripe = slide.shapes.add_shape(
        1, Inches(0), Inches(1.1), Inches(0.08), Inches(6.4)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = hex_to_rgb(palette["accent"])
    stripe.line.fill.background()

    # Body text / bullets
    content_top = Inches(1.3)
    content_left = Inches(0.5)
    content_width = Inches(11.8)
    content_height = Inches(5.7)

    if slide_data["bullets"]:
        txb2 = slide.shapes.add_textbox(
            content_left, content_top, content_width, content_height
        )
        tf2 = txb2.text_frame
        tf2.word_wrap = True

        for i, bullet in enumerate(slide_data["bullets"]):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.level = bullet["level"]
            indent = Pt(bullet["level"] * 14)
            p.space_before = Pt(4 if bullet["level"] > 0 else 8)

            run = p.add_run()
            run.text = (
                "• " if bullet["level"] == 1 else "  ◦ " if bullet["level"] == 2 else ""
            ) + bullet["text"]
            run.font.size = Pt(15 if bullet["level"] > 0 else 17)
            run.font.bold = bullet["bold"]
            run.font.name = (
                palette["body_font"] if not bullet["bold"] else palette["header_font"]
            )
            run.font.color.rgb = hex_to_rgb(
                palette["dark"] if bullet["bold"] else palette["text"]
            )

    elif slide_data["body"].strip():
        txb2 = slide.shapes.add_textbox(
            content_left, content_top, content_width, content_height
        )
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        run = p.add_run()
        run.text = slide_data["body"].strip()
        run.font.size = Pt(16)
        run.font.color.rgb = hex_to_rgb(palette["text"])
        run.font.name = palette["body_font"]

    return slide


# ─── Palette Registry ────────────────────────────────────────────────────────

PALETTES = {
    "midnight": {
        "dark": "1E2761",
        "light": "FFFFFF",
        "light_alt": "F0F4FF",
        "accent": "4A90D9",
        "muted": "CADCFC",
        "text": "1E2761",
        "header_font": "Georgia",
        "body_font": "Calibri",
    },
    "forest": {
        "dark": "2C5F2D",
        "light": "FFFFFF",
        "light_alt": "F4F9F4",
        "accent": "97BC62",
        "muted": "C8DDB0",
        "text": "1A3A1B",
        "header_font": "Trebuchet MS",
        "body_font": "Calibri",
    },
    "coral": {
        "dark": "2F3C7E",
        "light": "FFFFFF",
        "light_alt": "FFF8F0",
        "accent": "F96167",
        "muted": "F9E795",
        "text": "2F3C7E",
        "header_font": "Cambria",
        "body_font": "Calibri",
    },
    "terracotta": {
        "dark": "B85042",
        "light": "F9F6F0",
        "light_alt": "FFFFFF",
        "accent": "A7BEAE",
        "muted": "E7E8D1",
        "text": "3A1A18",
        "header_font": "Palatino",
        "body_font": "Calibri",
    },
    "ocean": {
        "dark": "065A82",
        "light": "FFFFFF",
        "light_alt": "EEF6FB",
        "accent": "02C39A",
        "muted": "B3D9EC",
        "text": "02243A",
        "header_font": "Arial Black",
        "body_font": "Calibri",
    },
    "charcoal": {
        "dark": "36454F",
        "light": "F2F2F2",
        "light_alt": "FFFFFF",
        "accent": "02C39A",
        "muted": "AAAAAA",
        "text": "1A1A1A",
        "header_font": "Calibri",
        "body_font": "Calibri Light",
    },
    "cherry": {
        "dark": "990011",
        "light": "FCF6F5",
        "light_alt": "FFFFFF",
        "accent": "2F3C7E",
        "muted": "D9A0A6",
        "text": "3A0008",
        "header_font": "Georgia",
        "body_font": "Calibri",
    },
}


# ─── Tool Function ────────────────────────────────────────────────────────────


async def generate_presentation(
    markdown_content: str,
    output_path: str,
    palette_name: str = "midnight",
    return_base64: bool = False,
    business_name: Optional[str] = None,
) -> dict:
    """
    Tool function: converts Markdown content into a .pptx presentation.

    Args:
        markdown_content (str):  Full markdown string from the upstream agent.
        output_path      (str):  File path to save the .pptx. Pass None to skip saving.
        palette_name     (str):  One of: midnight | forest | coral | terracotta |
                                         ocean | charcoal | cherry
        return_base64    (bool): If True, also returns the file as a base64 string
                                 (useful when the tool runs in a sandboxed environment).

    Returns:
        {
          "success":    bool,
          "slide_count": int,
          "output_path": str | None,
          "base64":     str | None,   # only when return_base64=True
          "error":      str | None,
        }
    """
    try:
        print("=============> [generate_presentation]")

        palette = PALETTES.get(palette_name, PALETTES["midnight"])
        slides_data = await ai_parse_markdown_to_slides(markdown_content)
        slides_data = normalize_bullets(slides_data)

        if not slides_data:
            print("=============> [generate_presentation] No Slides Data")

            return {
                "success": False,
                "slide_count": 0,
                "output_path": None,
                "base64": None,
                "error": "No slides parsed from markdown.",
            }

        print("=============> [generate_presentation] Parsed Slides Data", slides_data)

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        content_index = 0
        for slide_data in slides_data:
            if slide_data["is_title_slide"]:
                build_title_slide(prs, slide_data, palette)
            else:
                build_content_slide(prs, slide_data, palette, content_index)
                content_index += 1

        # Save to file
        if output_path:
            prs.save(output_path)

        print("=============> [generate_presentation] Saved to file", output_path)

        # Optionally encode to base64
        b64 = None
        if return_base64:
            buffer = io.BytesIO()
            prs.save(buffer)
            b64 = base64.b64encode(buffer.getvalue()).decode("utf-8")

        print("=============> [generate_presentation] Returning result")

        return {
            "success": True,
            "slide_count": len(slides_data),
            "output_path": output_path,
            "base64": b64,
            "error": None,
        }

    except Exception as e:
        return {
            "success": False,
            "slide_count": 0,
            "output_path": None,
            "base64": None,
            "error": str(e),
        }
